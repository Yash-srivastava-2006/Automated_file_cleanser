import json
import re
import shutil
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional

import cv2
import fitz
import pandas as pd
import pdfplumber
import pytesseract
import spacy
from pptx import Presentation
from ultralytics import YOLO


EMAIL_PATTERN = re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b")
PHONE_PATTERN = re.compile(r"\b(?:\+?\d{1,3}[\s.-]?)?(?:\(?\d{3}\)?[\s.-]?)\d{3}[\s.-]?\d{4}\b")
AADHAAR_PATTERN = re.compile(r"\b\d{4}[\s-]?\d{4}[\s-]?\d{4}\b")
SSN_PATTERN = re.compile(r"\b\d{3}-?\d{2}-?\d{4}\b")

TAG_MAP = {
    "PERSON": "[REDACTED_NAME]",
    "EMAIL": "[REDACTED_EMAIL]",
    "PHONE": "[REDACTED_PHONE]",
    "AADHAAR": "[REDACTED_ID]",
    "SSN": "[REDACTED_ID]",
}

SUPPORTED_EXTENSIONS = {"pdf", "png", "jpg", "jpeg", "pptx", "xlsx", "txt"}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg"}


@dataclass
class ProcessedArtifact:
    file_name: str
    file_type: str
    file_description: str
    key_findings: List[str]
    entities: List[Dict[str, str]]
    raw_text: str
    cleaned_text: str
    masked_file_path: str
    raw_text_path: str
    clean_text_path: str
    json_path: str


class PIIDetector:
    def __init__(self) -> None:
        self.model_name = "en_core_web_trf"
        self.nlp = self._load_spacy_model()

    def _load_spacy_model(self):
        try:
            return spacy.load(self.model_name)
        except Exception:
            # Fallback keeps app usable when transformer model is missing.
            try:
                self.model_name = "en_core_web_sm"
                return spacy.load(self.model_name)
            except Exception:
                # Last fallback: blank pipeline so regex detections still work.
                self.model_name = "blank_en"
                return spacy.blank("en")

    def detect(self, text: str) -> List[Dict[str, str]]:
        entities: List[Dict[str, str]] = []

        doc = self.nlp(text)
        for ent in doc.ents:
            if ent.label_ == "PERSON" and ent.text.strip():
                entities.append(
                    {
                        "entity": "PERSON",
                        "text": ent.text,
                        "start": ent.start_char,
                        "end": ent.end_char,
                        "source": f"spacy:{self.model_name}",
                    }
                )

        entities.extend(self._regex_entities(text, EMAIL_PATTERN, "EMAIL"))
        entities.extend(self._regex_entities(text, PHONE_PATTERN, "PHONE"))
        entities.extend(self._regex_entities(text, AADHAAR_PATTERN, "AADHAAR"))
        entities.extend(self._regex_entities(text, SSN_PATTERN, "SSN"))

        return self._deduplicate_entities(entities)

    def _regex_entities(self, text: str, pattern: re.Pattern, label: str) -> List[Dict[str, str]]:
        output: List[Dict[str, str]] = []
        for match in pattern.finditer(text):
            output.append(
                {
                    "entity": label,
                    "text": match.group(0),
                    "start": match.start(),
                    "end": match.end(),
                    "source": "regex",
                }
            )
        return output

    def _deduplicate_entities(self, entities: List[Dict[str, str]]) -> List[Dict[str, str]]:
        entities = sorted(entities, key=lambda x: (x["start"], -(x["end"] - x["start"])))
        merged: List[Dict[str, str]] = []
        for item in entities:
            overlap = False
            for existing in merged:
                if not (item["end"] <= existing["start"] or item["start"] >= existing["end"]):
                    overlap = True
                    break
            if not overlap:
                merged.append(item)
        return merged

    def mask_text(self, text: str, entities: List[Dict[str, str]]) -> str:
        masked = text
        for ent in sorted(entities, key=lambda x: x["start"], reverse=True):
            tag = TAG_MAP.get(ent["entity"], f"[REDACTED_{ent['entity']}]")
            masked = masked[: ent["start"]] + tag + masked[ent["end"] :]
        return masked


class FileCleansingPipeline:
    def __init__(self, root_dir: str) -> None:
        self.root = Path(root_dir)
        self.upload_dir = self.root / "uploads"
        self.output_dir = self.root / "outputs"
        self.masked_dir = self.output_dir / "masked"
        self.raw_dir = self.output_dir / "raw_text"
        self.clean_dir = self.output_dir / "clean_text"
        self.json_dir = self.output_dir / "json"

        for directory in [
            self.upload_dir,
            self.output_dir,
            self.masked_dir,
            self.raw_dir,
            self.clean_dir,
            self.json_dir,
        ]:
            directory.mkdir(parents=True, exist_ok=True)

        self.detector = PIIDetector()
        self.face_model = self._load_face_model()

    def _load_face_model(self):
        weight_path = self.root / "yolov8n-face.pt"
        if not weight_path.exists():
            return None
        try:
            return YOLO(str(weight_path))
        except Exception:
            return None

    def detect_file_type(self, file_path: Path) -> str:
        return file_path.suffix.lower().replace(".", "")

    def process_file(self, file_path: Path) -> ProcessedArtifact:
        file_type = self.detect_file_type(file_path)
        raw_text = self.extract_text(file_path, file_type)
        normalized_text = self.normalize_text(raw_text)

        entities = self.detector.detect(normalized_text)
        cleaned_text = self.detector.mask_text(normalized_text, entities)

        masked_file_path = self.mask_file(file_path, file_type, entities, cleaned_text)

        raw_text_path = self.raw_dir / f"{file_path.stem}_raw.txt"
        clean_text_path = self.clean_dir / f"{file_path.stem}_clean.txt"
        raw_text_path.write_text(raw_text, encoding="utf-8")
        clean_text_path.write_text(cleaned_text, encoding="utf-8")

        file_description = self.generate_description(cleaned_text)
        key_findings = self.generate_key_findings(cleaned_text)

        json_payload = {
            "file_name": file_path.name,
            "file_type": file_type,
            "file_description": file_description,
            "key_findings": key_findings,
        }
        json_path = self.json_dir / f"{file_path.stem}_report.json"
        json_path.write_text(json.dumps(json_payload, indent=2), encoding="utf-8")

        return ProcessedArtifact(
            file_name=file_path.name,
            file_type=file_type,
            file_description=file_description,
            key_findings=key_findings,
            entities=entities,
            raw_text=raw_text,
            cleaned_text=cleaned_text,
            masked_file_path=str(masked_file_path),
            raw_text_path=str(raw_text_path),
            clean_text_path=str(clean_text_path),
            json_path=str(json_path),
        )

    def extract_text(self, file_path: Path, file_type: str) -> str:
        if file_type in IMAGE_EXTENSIONS:
            image = cv2.imread(str(file_path))
            return pytesseract.image_to_string(image) if image is not None else ""
        if file_type == "pdf":
            return self._extract_pdf_text(file_path)
        if file_type == "pptx":
            return self._extract_pptx_text(file_path)
        if file_type == "xlsx":
            return self._extract_xlsx_text(file_path)
        if file_type == "txt":
            return file_path.read_text(encoding="utf-8", errors="ignore")
        return ""

    def _extract_pdf_text(self, file_path: Path) -> str:
        text_parts: List[str] = []
        try:
            with pdfplumber.open(str(file_path)) as pdf:
                for page in pdf.pages:
                    text_parts.append(page.extract_text() or "")
        except Exception:
            text_parts = []

        joined = "\n".join(text_parts).strip()
        if joined:
            return joined

        # Fallback to PyMuPDF when pdfplumber has no readable layer.
        doc = fitz.open(str(file_path))
        fallback_parts = [page.get_text("text") for page in doc]
        doc.close()
        return "\n".join(fallback_parts)

    def _extract_pptx_text(self, file_path: Path) -> str:
        prs = Presentation(str(file_path))
        lines: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text:
                    lines.append(shape.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                lines.append(cell.text)
        return "\n".join(lines)

    def _extract_xlsx_text(self, file_path: Path) -> str:
        sheets = pd.read_excel(file_path, sheet_name=None, dtype=str)
        chunks: List[str] = []
        for sheet_name, df in sheets.items():
            chunks.append(f"Sheet: {sheet_name}")
            chunks.append(df.fillna("").to_string(index=False))
        return "\n\n".join(chunks)

    def normalize_text(self, text: str) -> str:
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def mask_file(self, file_path: Path, file_type: str, entities: List[Dict[str, str]], cleaned_text: str) -> Path:
        output_name = f"masked_{uuid.uuid4()}_{file_path.name}"
        output_path = self.masked_dir / output_name

        if file_type in IMAGE_EXTENSIONS:
            self._mask_image(file_path, output_path, entities)
            return output_path
        if file_type == "pdf":
            self._mask_pdf(file_path, output_path, entities)
            return output_path
        if file_type == "pptx":
            self._mask_pptx(file_path, output_path)
            return output_path
        if file_type == "xlsx":
            self._mask_xlsx(file_path, output_path)
            return output_path
        if file_type == "txt":
            output_path.write_text(cleaned_text, encoding="utf-8")
            return output_path

        shutil.copy2(file_path, output_path)
        return output_path

    def _mask_image(self, file_path: Path, output_path: Path, entities: List[Dict[str, str]]) -> None:
        image = cv2.imread(str(file_path))
        if image is None:
            shutil.copy2(file_path, output_path)
            return

        if self.face_model is not None:
            try:
                results = self.face_model.predict(source=image, verbose=False)
                for result in results:
                    if result.boxes is None:
                        continue
                    for box in result.boxes.xyxy.cpu().numpy().astype(int):
                        x1, y1, x2, y2 = box.tolist()
                        roi = image[y1:y2, x1:x2]
                        if roi.size > 0:
                            image[y1:y2, x1:x2] = cv2.GaussianBlur(roi, (51, 51), 30)
            except Exception:
                pass

        entity_tokens = set()
        for ent in entities:
            for token in re.split(r"\W+", ent["text"].lower()):
                if token:
                    entity_tokens.add(token)

        ocr_data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
        for idx, token in enumerate(ocr_data.get("text", [])):
            word = token.strip()
            if not word:
                continue
            clean_word = re.sub(r"\W+", "", word.lower())
            should_mask = (
                clean_word in entity_tokens
                or bool(EMAIL_PATTERN.search(word))
                or bool(PHONE_PATTERN.search(word))
                or bool(AADHAAR_PATTERN.search(word))
                or bool(SSN_PATTERN.search(word))
            )
            if should_mask:
                x = int(ocr_data["left"][idx])
                y = int(ocr_data["top"][idx])
                w = int(ocr_data["width"][idx])
                h = int(ocr_data["height"][idx])
                cv2.rectangle(image, (x, y), (x + w, y + h), (0, 0, 0), thickness=-1)

        cv2.imwrite(str(output_path), image)

    def _mask_pdf(self, file_path: Path, output_path: Path, entities: List[Dict[str, str]]) -> None:
        doc = fitz.open(str(file_path))
        for page in doc:
            for ent in entities:
                found = page.search_for(ent["text"])
                for area in found:
                    page.add_redact_annot(area, fill=(0, 0, 0))
            page.apply_redactions()
        doc.save(str(output_path))
        doc.close()

    def _mask_pptx(self, file_path: Path, output_path: Path) -> None:
        prs = Presentation(str(file_path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text:
                    masked_text = self.detector.mask_text(shape.text, self.detector.detect(shape.text))
                    shape.text = masked_text
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            cell_text = cell.text or ""
                            masked_cell = self.detector.mask_text(cell_text, self.detector.detect(cell_text))
                            cell.text = masked_cell
        prs.save(str(output_path))

    def _mask_xlsx(self, file_path: Path, output_path: Path) -> None:
        sheets = pd.read_excel(file_path, sheet_name=None, dtype=str)
        with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
            for name, df in sheets.items():
                masked_df = df.fillna("").applymap(
                    lambda value: self.detector.mask_text(str(value), self.detector.detect(str(value)))
                    if str(value).strip()
                    else value
                )
                masked_df.to_excel(writer, index=False, sheet_name=name)

    def generate_description(self, cleaned_text: str) -> str:
        if not cleaned_text.strip():
            return "No readable text found after extraction."

        sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", cleaned_text) if s.strip()]
        if not sentences:
            return "Document content was extracted and cleansed for analysis."
        preview = " ".join(sentences[:2])
        return preview[:300]

    def generate_key_findings(self, cleaned_text: str) -> List[str]:
        text = cleaned_text.lower()
        findings: List[str] = []

        rules = [
            (["biometric", "fingerprint", "face recognition"], "Uses biometric authentication controls."),
            (["manual", "paper", "spreadsheet"], "Manual processes may introduce operational errors."),
            (["timestamp", "audit log", "access log"], "Tracks access events with time-based logging."),
            (["encryption", "aes", "tls", "ssl"], "Applies encryption mechanisms for data protection."),
            (["mfa", "multi-factor", "2fa"], "Enforces multi-factor authentication for account security."),
            (["firewall", "allow", "deny", "port"], "Defines explicit network access and firewall controls."),
            (["policy", "compliance", "governance"], "References compliance or governance policy controls."),
            (["backup", "recovery", "disaster"], "Mentions resilience through backup or recovery planning."),
        ]

        for keywords, finding in rules:
            if any(keyword in text for keyword in keywords):
                findings.append(finding)

        if len(findings) < 3:
            sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", cleaned_text) if len(s.strip()) > 25]
            for sentence in sentences:
                normalized = sentence.replace("\n", " ")
                findings.append(f"Highlights: {normalized[:140]}")
                if len(findings) >= 5:
                    break

        if not findings:
            findings = ["No significant themes were detected from extracted content."]

        while len(findings) < 3:
            findings.append("Additional review recommended to derive more domain-specific insights.")

        return findings[:5]


def safe_filename(file_name: str) -> str:
    return re.sub(r"[^A-Za-z0-9._-]", "_", file_name)


def save_uploaded_file(uploaded_file, upload_dir: Path) -> Optional[Path]:
    extension = uploaded_file.name.rsplit(".", 1)[-1].lower() if "." in uploaded_file.name else ""
    if extension not in SUPPORTED_EXTENSIONS:
        return None

    file_name = f"{uuid.uuid4()}_{safe_filename(uploaded_file.name)}"
    destination = upload_dir / file_name
    destination.write_bytes(uploaded_file.getbuffer())
    return destination
